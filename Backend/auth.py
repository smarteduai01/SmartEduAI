import os
import json
import tempfile
import requests
from dotenv import load_dotenv
from pptx import Presentation
import docx2txt
import PyPDF2
from langchain.text_splitter import RecursiveCharacterTextSplitter
from flask import Flask, request, jsonify, session
from flask_cors import CORS
from pymongo.mongo_client import MongoClient
import bcrypt
from functools import wraps
from datetime import datetime
from urllib.parse import quote_plus
from pymongo.mongo_client import MongoClient

# ------------------- FLASK SETUP -------------------
app = Flask(__name__)
CORS(app, origins=["http://localhost:5173"], supports_credentials=True)  # React frontend
app.secret_key = os.getenv("SECRET_KEY", "supersecretkey")

# Load environment variables
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
MONGO_URI = os.getenv("MONGO_URI")

# MongoDB Connection
username = quote_plus("JayanthSrinivas02")
password = quote_plus("gFVKHiFFX86oZ6wj")
uri = f"mongodb+srv://{username}:{password}@smarteduai.bo81fvz.mongodb.net/?appName=SmartEduAI"

client = MongoClient(uri)
db = client["QuizAI"]
users_col = db["users"]
quiz_results_col = db["quiz_results"]

# Gemini endpoint
GEMINI_ENDPOINT = (
    f"https://generativelanguage.googleapis.com/v1beta/models/"
    f"gemini-2.5-flash:generateContent?key={GEMINI_API_KEY}"
)

# ------------------- AUTH HELPERS -------------------
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'username' not in session:
            return jsonify({"error": "Authentication required"}), 401
        return f(*args, **kwargs)
    return decorated_function

# ------------------- AUTH ROUTES -------------------
@app.route("/signup", methods=["POST"])
def signup():
    data = request.get_json()
    username = data.get("username")
    password = data.get("password")

    if not username or not password:
        return jsonify({"error": "Username and password required"}), 400

    if users_col.find_one({'username': username}):
        return jsonify({"error": "Username already exists"}), 400

    hashed_pw = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())
    users_col.insert_one({'username': username, 'password': hashed_pw})

    return jsonify({"message": "User created successfully"}), 201


@app.route("/login", methods=["POST"])
def login():
    data = request.get_json()
    username = data.get("username")
    password = data.get("password")

    if not username or not password:
        return jsonify({"error": "Username and password required"}), 400

    user = users_col.find_one({'username': username})
    if user and bcrypt.checkpw(password.encode('utf-8'), user['password']):
        session['username'] = username
        return jsonify({"message": "Login successful", "username": username}), 200
    else:
        return jsonify({"error": "Invalid username or password"}), 401


@app.route("/logout", methods=["POST"])
def logout():
    session.clear()
    return jsonify({"message": "Logged out successfully"}), 200


@app.route("/check_auth", methods=["GET"])
def check_auth():
    if 'username' in session:
        return jsonify({"authenticated": True, "username": session['username']}), 200
    else:
        return jsonify({"authenticated": False}), 200

# ------------------- MCQ GENERATION -------------------
def extract_text(file_path, ext):
    text = ""
    if ext == "pdf":
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                if page_text := page.extract_text():
                    text += page_text + "\n"
    elif ext == "docx":
        text = docx2txt.process(file_path)
    elif ext == "pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text.strip()


def call_gemini(prompt):
    headers = {"Content-Type": "application/json"}
    data = {"contents": [{"parts": [{"text": prompt}]}]}
    resp = requests.post(GEMINI_ENDPOINT, headers=headers, json=data)
    resp.raise_for_status()
    json_response = resp.json()
    return json_response["candidates"][0]["content"]["parts"][0]["text"]


@app.route("/generate_mcq", methods=["POST"])
@login_required
def generate_mcq():
    file = request.files.get("file")
    if not file:
        return jsonify({"error": "No file uploaded"}), 400

    ext = file.filename.split(".")[-1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix="."+ext) as tmp:
        file.save(tmp.name)
        file_path = tmp.name

    text = extract_text(file_path, ext)
    os.unlink(file_path)

    if not text:
        return jsonify({"error": "No readable text found in file"}), 400

    splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
    chunks = splitter.split_text(text)
    main_context = " ".join(chunks[:3])

    prompt = f"""
You are an expert educational AI system. Generate exactly 10 multiple-choice questions (MCQs)
from the given context.

### Context:
{main_context}

### Instructions:
- Generate exactly 10 questions
- Each question must have 4 options and one correct answer.
- Include difficulty level (Easy, Medium, Hard).
- Return strictly JSON format:

[
  {{
    "question": "Question text?",
    "options": ["A", "B", "C", "D"],
    "correct_answer": "A",
    "difficulty": "Easy"
  }}
]
"""
    try:
        response_text = call_gemini(prompt)
        if response_text.strip().startswith("```json"):
            response_text = response_text.strip()[7:-3].strip()
        mcq_json = json.loads(response_text)

        session['current_quiz'] = {
            'questions': mcq_json,
            'filename': file.filename,
            'started_at': datetime.now().isoformat()
        }

        return jsonify(mcq_json)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ------------------- QUIZ SUBMISSION -------------------
@app.route("/submit_quiz", methods=["POST"])
@login_required
def submit_quiz():
    data = request.get_json()
    answers = data.get("answers")

    if 'current_quiz' not in session:
        return jsonify({"error": "No active quiz session"}), 400

    quiz_data = session['current_quiz']
    questions = quiz_data['questions']

    correct_count = 0
    total = len(questions)
    results = []

    for i, q in enumerate(questions):
        user_answer = answers[i] if i < len(answers) else None
        is_correct = user_answer == q['correct_answer']
        if is_correct:
            correct_count += 1
        results.append({
            'question': q['question'],
            'user_answer': user_answer,
            'correct_answer': q['correct_answer'],
            'is_correct': is_correct,
            'difficulty': q['difficulty']
        })

    percentage = (correct_count / total) * 100

    quiz_results_col.insert_one({
        'username': session['username'],
        'filename': quiz_data['filename'],
        'started_at': quiz_data['started_at'],
        'completed_at': datetime.now().isoformat(),
        'total_questions': total,
        'correct_answers': correct_count,
        'score_percentage': percentage,
        'detailed_results': results
    })

    session.pop('current_quiz', None)

    return jsonify({
        'score': correct_count,
        'total': total,
        'percentage': percentage,
        'results': results
    }), 200

# ------------------- QUIZ HISTORY -------------------
@app.route("/quiz_history", methods=["GET"])
@login_required
def quiz_history():
    username = session['username']
    history = list(quiz_results_col.find({'username': username}, {'_id': 0}).sort('completed_at', -1))
    return jsonify(history), 200


if __name__ == "__main__":
    app.run(port=5000, debug=True)
